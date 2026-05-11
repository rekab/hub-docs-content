"""Stage C — Cascade classify. One model, three input modes (text / +image@280 / +image@560)."""
from __future__ import annotations

import hashlib
import json
import os
import re
import subprocess
import sys
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from . import _common

if False:  # pragma: no cover — for type checkers only
    from .llm import LMStudioClient

MANIFEST_PATH = _common.REPO_ROOT / "manifest.toml"

CONFIDENCE_CAPPED = "text_only_capped"
TEXT_TRUNCATE_CHARS = 24_000  # ~6k tokens worth


# ----- Kind detection --------------------------------------------------------

DOC_EXTS = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".txt", ".rtf", ".odt", ".odp"}
SPREADSHEET_EXTS = {".xlsx", ".xls", ".csv", ".ods"}
IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".heic", ".gif", ".webp", ".svg"}

SLIDE_DECK_MIMES = {
    "application/vnd.google-apps.presentation",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
}


def _kind_of(saved_path: Path, mime: str) -> str:
    ext = saved_path.suffix.lower()
    if mime.startswith("image/") or ext in IMAGE_EXTS:
        return "image"
    if ext in SPREADSHEET_EXTS:
        return "spreadsheet"
    return "document"


# ----- Manifest categories (anchor) ------------------------------------------


def _load_manifest_categories() -> list[str]:
    if not MANIFEST_PATH.exists():
        return []
    cats: list[str] = []
    seen: set[str] = set()
    for line in MANIFEST_PATH.read_text().splitlines():
        m = re.match(r'\s*category\s*=\s*"([^"]+)"', line)
        if m and m.group(1) not in seen:
            cats.append(m.group(1))
            seen.add(m.group(1))
    return cats


# ----- Text extraction -------------------------------------------------------


def _extract_pdf_text(path: Path) -> str:
    import fitz  # pymupdf
    with fitz.open(path) as doc:
        return "\n".join(page.get_text() for page in doc)


def _extract_docx_text(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    parts: list[str] = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            parts.extend(cell.text for cell in row.cells)
    return "\n".join(parts)


def _extract_pptx_text(path: Path) -> str:
    from pptx import Presentation
    pres = Presentation(str(path))
    parts: list[str] = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    parts.append(para.text)
    return "\n".join(parts)


def _extract_xlsx_text(path: Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(path), read_only=True, data_only=True)
    sheet = wb[wb.sheetnames[0]]
    rows: list[str] = []
    for r, row in enumerate(sheet.iter_rows(values_only=True)):
        if r >= 50:
            break
        cells = [str(c) if c is not None else "" for c in row[:50]]
        rows.append(" | ".join(cells))
    return "\n".join(rows)


def _extract_csv_text(path: Path) -> str:
    import csv as csvmod
    out: list[str] = []
    with path.open("r", encoding="utf-8", errors="replace", newline="") as f:
        reader = csvmod.reader(f)
        for r, row in enumerate(reader):
            if r >= 50:
                break
            out.append(" | ".join(row[:50]))
    return "\n".join(out)


def _convert_to_pdf_via_libreoffice(src: Path) -> Path:
    """Returns path to a freshly-rendered .pdf in a temp dir owned by caller."""
    outdir = Path(tempfile.mkdtemp(prefix="hubmine-soffice-"))
    soffice = os.environ.get("SOFFICE_BIN", "soffice")
    proc = subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(outdir), str(src)],
        capture_output=True,
        text=True,
        timeout=180,
    )
    if proc.returncode != 0:
        raise RuntimeError(f"libreoffice failed for {src.name}: {proc.stderr.strip()[:400]}")
    pdf = outdir / (src.stem + ".pdf")
    if not pdf.exists():
        raise RuntimeError(f"libreoffice produced no pdf for {src.name} (in {outdir})")
    return pdf


def _extract_text(saved_path: Path) -> str:
    ext = saved_path.suffix.lower()
    if ext == ".pdf":
        return _extract_pdf_text(saved_path)
    if ext == ".docx":
        return _extract_docx_text(saved_path)
    if ext == ".pptx":
        return _extract_pptx_text(saved_path)
    if ext == ".xlsx":
        return _extract_xlsx_text(saved_path)
    if ext == ".csv":
        return _extract_csv_text(saved_path)
    if ext in {".doc", ".ppt", ".rtf", ".odt", ".odp"}:
        pdf = _convert_to_pdf_via_libreoffice(saved_path)
        return _extract_pdf_text(pdf)
    if ext == ".txt":
        return saved_path.read_text(encoding="utf-8", errors="replace")
    return ""


# ----- Image rendering for vision passes -------------------------------------


def _render_pages_for_vision(saved_path: Path, mime: str, *, max_pages: int = 2, dpi: int = 144) -> list[Path]:
    """Render the first up-to-N pages of a document to PNGs in a temp dir.
    For genuine image files, returns [saved_path] directly (vision-mode passes raw bytes).
    """
    ext = saved_path.suffix.lower()
    if mime.startswith("image/") or ext in IMAGE_EXTS:
        return [saved_path]

    if ext == ".pdf":
        pdf = saved_path
    elif ext in {".docx", ".doc", ".pptx", ".ppt", ".rtf", ".odt", ".odp", ".xlsx", ".xls"}:
        pdf = _convert_to_pdf_via_libreoffice(saved_path)
    else:
        return []

    import fitz  # pymupdf
    out_dir = Path(tempfile.mkdtemp(prefix="hubmine-pages-"))
    out: list[Path] = []
    with fitz.open(pdf) as doc:
        for i, page in enumerate(doc):
            if i >= max_pages:
                break
            pix = page.get_pixmap(dpi=dpi)
            target = out_dir / f"page{i + 1}.png"
            pix.save(str(target))
            out.append(target)
    return out


# ----- LLM prompts -----------------------------------------------------------

CLASSIFY_SYSTEM_TEMPLATE = """You are classifying a document mined from the
Seattle Emergency Hubs community Google Drive, deciding whether it belongs
in the Hub Reference Library (curated PDF handouts served via captive
portal at hub sites during emergencies).

Existing manifest categories (anchors — prefer these; only propose a new
category if the document genuinely doesn't fit):
{categories}

Scope rules:
  - KEEP English citizen-facing emergency handouts AND hub-captain
    operational docs (org charts, contact templates, run-of-show, USSF
    talking points).
  - SKIP translations of any kind.
  - SKIP internal admin: rosters, contact lists, supplies tracking,
    distribution logs, meeting minutes, board agendas.

Audience taxonomy: one of "citizen", "hub_captain", "both".

Reply with a single JSON object exactly matching this schema:
{{
  "keep": <bool>,
  "slug": "<lowercase-hyphenated-filename-without-extension>",
  "title": "<title-case display title, < 80 chars>",
  "description": "<one or two sentences for the operator>",
  "category": "<existing category, or new short one>",
  "lang": "<two-letter ISO code, e.g. en, es>",
  "audience": "<citizen|hub_captain|both>",
  "useful_for_phone": <bool>,
  "confidence": <0.0..1.0>
}}
""".strip()


def _build_system_prompt() -> str:
    cats = _load_manifest_categories()
    cat_lines = "\n".join(f"  - {c}" for c in cats) if cats else "  (none defined yet)"
    return CLASSIFY_SYSTEM_TEMPLATE.format(categories=cat_lines)


def _user_prompt(*, sidecar: dict[str, Any], text: str, kind: str) -> str:
    parts = [
        f"Filename: {sidecar.get('original_name', '')}",
        f"Folder path: {sidecar.get('original_path', '')}",
        f"Owner: {sidecar.get('owner', '')}",
        f"MimeType: {sidecar.get('mimeType', '')}",
        f"File kind: {kind}",
    ]
    if text:
        snippet = text[:TEXT_TRUNCATE_CHARS]
        parts.append("")
        parts.append("Extracted content (may be truncated):")
        parts.append(snippet)
    else:
        parts.append("")
        parts.append("(No extracted text — judge from the image and filename.)")
    return "\n".join(parts)


def _prompt_sha(system: str, user_template_for_kind: str) -> str:
    h = hashlib.sha256()
    h.update(system.encode("utf-8"))
    h.update(b"\x1e")
    h.update(user_template_for_kind.encode("utf-8"))
    return h.hexdigest()[:16]


# ----- Cache helpers ---------------------------------------------------------


def _classified_path(file_id: str) -> Path:
    return _common.CLASSIFIED / f"{file_id}.meta.json"


def _extracted_path(file_id: str) -> Path:
    return _common.EXTRACTED / f"{file_id}.text.json"


def _confidence_below(value: Any, threshold: float) -> bool:
    if isinstance(value, (int, float)):
        return float(value) < threshold
    return False


# ----- Per-file classification ----------------------------------------------


@dataclass
class ClassifyResult:
    file_id: str
    status: str  # 'classified' | 'cached' | 'error'
    pass_taken: int  # 0 (image: skipped pass1), 1/2/3 (highest pass)
    note: str = ""


def _load_sidecar(file_id: str) -> dict[str, Any]:
    p = _common.DOWNLOADS / file_id / "_meta.json"
    if not p.exists():
        raise FileNotFoundError(f"no sidecar for fileId={file_id}")
    return json.loads(p.read_text())


def _saved_path(sidecar: dict[str, Any]) -> Path:
    return _common.DOWNLOADS / sidecar["file_id"] / sidecar["saved_filename"]


def _ensure_extracted(sidecar: dict[str, Any], saved_path: Path) -> str:
    cache_path = _extracted_path(sidecar["file_id"])
    if cache_path.exists():
        try:
            cached = json.loads(cache_path.read_text())
            if cached.get("content_sha256") == sidecar["bytes_sha256"]:
                return cached.get("text", "")
        except Exception:
            pass
    text = _extract_text(saved_path)
    _common.atomic_write_json(cache_path, {
        "file_id": sidecar["file_id"],
        "content_sha256": sidecar["bytes_sha256"],
        "text": text,
        "text_chars": len(text),
    })
    return text


def _cache_is_fresh(
    cached: dict[str, Any],
    *,
    content_sha256: str,
    model_id: str,
    prompt_sha: str,
    visual_token_budget: int,
) -> bool:
    keys = cached.get("cache_key", {})
    return (
        keys.get("content_sha256") == content_sha256
        and keys.get("model_id") == model_id
        and keys.get("prompt_sha") == prompt_sha
        and int(keys.get("visual_token_budget", -1)) == visual_token_budget
    )


def classify_file(client: "LMStudioClient", file_id: str) -> ClassifyResult:
    sidecar = _load_sidecar(file_id)
    saved_path = _saved_path(sidecar)
    if not saved_path.exists():
        return ClassifyResult(file_id, "error", 0, "saved file missing")

    mime = sidecar.get("mimeType", "")
    kind = _kind_of(saved_path, mime)
    system = _build_system_prompt()
    prompt_sha = _prompt_sha(system, kind)
    out_path = _classified_path(file_id)

    # Kind B — spreadsheets: text-only Pass 1, hard skip vision.
    if kind == "spreadsheet":
        existing = _read_cached(out_path)
        if existing and _cache_is_fresh(
            existing,
            content_sha256=sidecar["bytes_sha256"],
            model_id=client.model,
            prompt_sha=prompt_sha,
            visual_token_budget=0,
        ):
            return ClassifyResult(file_id, "cached", 1)
        text = _ensure_extracted(sidecar, saved_path)
        parsed = _call_llm(client, system, sidecar, text, kind, visual_token_budget=0)
        if _confidence_below(parsed.get("confidence"), 0.7):
            parsed["confidence"] = CONFIDENCE_CAPPED
            parsed["keep"] = False
        _write_classified(out_path, sidecar, parsed, pass_taken=1, kind=kind, model=client.model,
                          prompt_sha=prompt_sha, visual_token_budget=0)
        return ClassifyResult(file_id, "classified", 1)

    # Kind C — images: skip text, skip Pass 1, go straight to Pass 2 / 3.
    if kind == "image":
        existing = _read_cached(out_path)
        # If a cache entry exists at budget 560 already, accept it.
        for budget in (560, 280):
            if existing and _cache_is_fresh(
                existing,
                content_sha256=sidecar["bytes_sha256"],
                model_id=client.model,
                prompt_sha=prompt_sha,
                visual_token_budget=budget,
            ):
                return ClassifyResult(file_id, "cached", 2 if budget == 280 else 3)
        # Pass 2 with the image bytes directly from disk.
        parsed = _call_llm(client, system, sidecar, text="", kind=kind,
                           images=[saved_path], visual_token_budget=280)
        pass_taken = 2
        if _confidence_below(parsed.get("confidence"), 0.7):
            parsed = _call_llm(client, system, sidecar, text="", kind=kind,
                               images=[saved_path], visual_token_budget=560)
            pass_taken = 3
        _write_classified(out_path, sidecar, parsed, pass_taken=pass_taken, kind=kind, model=client.model,
                          prompt_sha=prompt_sha,
                          visual_token_budget=280 if pass_taken == 2 else 560)
        return ClassifyResult(file_id, "classified", pass_taken)

    # Kind A — documents: Pass 1 text / Pass 2 +image@280 / Pass 3 +image@560.
    existing = _read_cached(out_path)
    for budget in (560, 280, 0):
        if existing and _cache_is_fresh(
            existing,
            content_sha256=sidecar["bytes_sha256"],
            model_id=client.model,
            prompt_sha=prompt_sha,
            visual_token_budget=budget,
        ):
            return ClassifyResult(file_id, "cached", {0: 1, 280: 2, 560: 3}[budget])

    text = _ensure_extracted(sidecar, saved_path)
    text_len = len(text.strip())

    # Pass 1: text-only.
    parsed = _call_llm(client, system, sidecar, text, kind, visual_token_budget=0)
    pass_taken = 1
    visual_token_budget = 0

    needs_pass2 = (
        _confidence_below(parsed.get("confidence"), 0.7)
        or (text_len < 100 and _looks_multipage(saved_path, mime))
        or _is_slide_deck(saved_path, mime)
    )
    if needs_pass2:
        try:
            pages = _render_pages_for_vision(saved_path, mime, max_pages=2, dpi=144)
            if pages:
                parsed = _call_llm(client, system, sidecar, text, kind,
                                   images=pages, visual_token_budget=280)
                pass_taken = 2
                visual_token_budget = 280

                needs_pass3 = (
                    _confidence_below(parsed.get("confidence"), 0.7)
                    and _looks_ocr_heavy(saved_path, text)
                )
                if needs_pass3:
                    parsed = _call_llm(client, system, sidecar, text, kind,
                                       images=pages, visual_token_budget=560)
                    pass_taken = 3
                    visual_token_budget = 560
        except Exception as e:
            print(f"[classify] vision pass failed for {file_id}: {e}", file=sys.stderr)

    _write_classified(out_path, sidecar, parsed, pass_taken=pass_taken, kind=kind, model=client.model,
                      prompt_sha=prompt_sha, visual_token_budget=visual_token_budget)
    return ClassifyResult(file_id, "classified", pass_taken)


def _looks_multipage(saved_path: Path, mime: str) -> bool:
    ext = saved_path.suffix.lower()
    if ext == ".pdf":
        try:
            import fitz
            with fitz.open(saved_path) as doc:
                return len(doc) >= 2
        except Exception:
            return False
    return ext in {".pptx", ".ppt", ".docx", ".doc"} or mime in SLIDE_DECK_MIMES


def _is_slide_deck(saved_path: Path, mime: str) -> bool:
    if mime in SLIDE_DECK_MIMES:
        return True
    return saved_path.suffix.lower() in {".pptx", ".ppt"}


def _looks_ocr_heavy(saved_path: Path, text: str) -> bool:
    size = saved_path.stat().st_size
    first_page_text = text[:500].strip()
    if size > 200_000 and len(first_page_text) < 50:
        return True
    return False


def _read_cached(out_path: Path) -> dict[str, Any] | None:
    if not out_path.exists():
        return None
    try:
        return json.loads(out_path.read_text())
    except Exception:
        return None


def _write_classified(
    out_path: Path,
    sidecar: dict[str, Any],
    parsed: dict[str, Any],
    *,
    pass_taken: int,
    kind: str,
    model: str,
    prompt_sha: str,
    visual_token_budget: int,
) -> None:
    # Normalize slug
    slug_in = parsed.get("slug") or _common.slugify(Path(sidecar.get("original_name", "untitled")).stem)
    parsed["slug"] = _common.slugify(slug_in)
    parsed.setdefault("title", sidecar.get("original_name", ""))
    parsed.setdefault("lang", "en")
    parsed.setdefault("audience", "citizen")
    parsed.setdefault("useful_for_phone", False)
    parsed.setdefault("category", "Uncategorized")
    parsed.setdefault("description", "")
    parsed.setdefault("keep", False)
    record = {
        "file_id": sidecar["file_id"],
        "kind": kind,
        "saved_filename": sidecar["saved_filename"],
        "result": parsed,
        "pass_taken": pass_taken,
        "cache_key": {
            "content_sha256": sidecar["bytes_sha256"],
            "model_id": model,
            "prompt_sha": prompt_sha,
            "visual_token_budget": visual_token_budget,
        },
    }
    _common.atomic_write_json(out_path, record)


def _call_llm(
    client: "LMStudioClient",
    system: str,
    sidecar: dict[str, Any],
    text: str,
    kind: str,
    *,
    images: list[Path] | None = None,
    visual_token_budget: int = 0,
) -> dict[str, Any]:
    user = _user_prompt(sidecar=sidecar, text=text, kind=kind)
    parsed, _raw = client.classify_json(
        system=system,
        user=user,
        images=images,
        visual_token_budget=visual_token_budget,
        max_tokens=1024,
    )
    return parsed


# ----- Driver ----------------------------------------------------------------


def _iter_downloaded_file_ids() -> list[str]:
    if not _common.DOWNLOADS.exists():
        return []
    return sorted(d.name for d in _common.DOWNLOADS.iterdir() if (d / "_meta.json").exists())


def run(*, limit: int | None = None) -> None:
    from .llm import LMStudioClient  # lazy: avoid forcing httpx/openai at import time

    _common.load_env()
    _common.ensure_dirs()

    client = LMStudioClient.from_env()
    print(f"[classify] probing LM Studio at {client.url} (model={client.model})", file=sys.stderr)
    client.probe()

    ids = _iter_downloaded_file_ids()
    if limit is not None:
        ids = ids[:limit]
    print(f"[classify] {len(ids)} downloaded files to classify", file=sys.stderr)

    stop = _common.StopRequested()
    stop.install()

    counters: dict[str, int] = {}
    for i, fid in enumerate(ids, 1):
        if stop.flagged:
            print("[classify] stop requested; exiting clean", file=sys.stderr)
            break
        try:
            res = classify_file(client, fid)
        except Exception as e:
            print(f"[classify] {fid}: error {e!r}", file=sys.stderr)
            counters["error"] = counters.get("error", 0) + 1
            continue
        counters[res.status] = counters.get(res.status, 0) + 1
        if res.status != "cached":
            print(f"[classify] {i}/{len(ids)} pass{res.pass_taken}: {fid}", file=sys.stderr)

    print(f"[classify] done. counters: {counters}", file=sys.stderr)
